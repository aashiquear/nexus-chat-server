"""
File Generator tool.

Generates PDF, DOCX, PPTX, and PNG files from content provided by the LLM.
Files are saved to the configured output directory and returned as downloadable
artifacts.

Supported formats
─────────────────
  pdf  – formatted text document (via fpdf2)
  doc  – Word document (via python-docx)
  ppt  – PowerPoint presentation (via python-pptx)
  png  – image from base64 data, text rendered to image (via Pillow),
         or AI-generated image from a text prompt (via Z-Image-Turbo)
"""

import asyncio
import base64
import json
import logging
import os
import re
import threading
import uuid
from pathlib import Path

from . import BaseTool, register_tool

logger = logging.getLogger(__name__)

# ─── Diffusion model singletons (lazy-loaded, thread-safe) ───
_pipe = None
_pipe_device = None
_pipe_lock = threading.Lock()

DIFFUSION_MODEL_ID = "Tongyi-MAI/Z-Image-Turbo"


def _ensure_diffusion_loaded(cache_dir: str | None = None):
    """Load the Z-Image-Turbo diffusion pipeline once (thread-safe).

    Mirrors the image_synthesizer pattern: module-level singleton with
    double-checked locking.  This function is **synchronous**; the caller
    (async execute()) must offload it to a thread so the event loop
    isn't frozen during the heavy from_pretrained() call.
    """
    global _pipe, _pipe_device

    if _pipe is not None:
        return

    with _pipe_lock:
        if _pipe is not None:
            return

        try:
            import torch
            from diffusers import DiffusionPipeline
        except ImportError as e:
            raise RuntimeError(
                "Text-to-image generation requires 'diffusers', 'transformers', "
                "'accelerate', and 'torch'. Install them with:\n"
                "  pip install diffusers transformers accelerate torch"
            ) from e

        _pipe_device = "cuda" if torch.cuda.is_available() else "cpu"
        logger.info("Loading %s on %s …", DIFFUSION_MODEL_ID, _pipe_device)

        load_kwargs = {}
        if cache_dir:
            load_kwargs["cache_dir"] = cache_dir

        torch_dtype = torch.bfloat16 if _pipe_device == "cuda" else torch.float32

        # Simple, direct load — exactly like image_synthesizer does with SmolVLM
        _pipe = DiffusionPipeline.from_pretrained(
            DIFFUSION_MODEL_ID,
            torch_dtype=torch_dtype,
            **load_kwargs,
        ).to(_pipe_device)

        logger.info("Diffusion model ready on %s", _pipe_device)


def _run_inference(prompt: str, title: str = ""):
    """Run the diffusion pipeline in a thread."""
    global _pipe, _pipe_device
    import torch

    full_prompt = f"{title}. {prompt}" if title else prompt
    full_prompt = full_prompt.strip()

    generator = None
    if _pipe_device == "cuda":
        generator = torch.Generator(_pipe_device).manual_seed(42)

    result = _pipe(
        prompt=full_prompt,
        height=1024,
        width=1024,
        num_inference_steps=9,
        guidance_scale=0.0,
        generator=generator,
    )
    return result.images[0]


@register_tool("file_generator")
class FileGeneratorTool(BaseTool):
    name = "file_generator"
    description = (
        "Generate PDF, DOCX, PPTX, or PNG files from content. "
        "Provide the file_type ('pdf', 'doc', 'ppt', 'png'), the content to embed, "
        "and optionally a title and filename. For PNG, content can be base64 image data, "
        "plain text rendered as an image, or a natural-language prompt for AI image generation. "
        "For PPT, separate slides with '---'. "
        "The tool creates the file, saves it, and returns a download link."
    )
    parameters = {
        "type": "object",
        "properties": {
            "file_type": {
                "type": "string",
                "enum": ["pdf", "doc", "ppt", "png"],
                "description": "Type of file to generate: pdf, doc, ppt, or png.",
            },
            "content": {
                "type": "string",
                "description": (
                    "Content to put in the file. For PDF/DOC/PPT: text or markdown. "
                    "For PNG: base64-encoded image data, text to render as an image, "
                    "or a natural-language prompt for AI image generation (e.g. 'a cat in a hat')."
                ),
            },
            "title": {
                "type": "string",
                "description": "Optional title for the document or image.",
            },
            "filename": {
                "type": "string",
                "description": "Optional suggested filename (including extension). If omitted, a UUID is generated.",
            },
        },
        "required": ["file_type", "content"],
    }

    async def execute(self, **kwargs) -> str:
        file_type = kwargs.get("file_type", "pdf").lower().strip().lstrip(".")
        content = kwargs.get("content", "")
        title = kwargs.get("title", "")
        suggested_filename = kwargs.get("filename", "")

        output_dir = Path(self.config.get("output_dir", "./data/downloads"))
        output_dir.mkdir(parents=True, exist_ok=True)

        if not content:
            return json.dumps({"error": "No content provided."})

        try:
            if file_type == "pdf":
                return self._generate_pdf(content, title, suggested_filename, output_dir)
            elif file_type == "doc":
                return self._generate_doc(content, title, suggested_filename, output_dir)
            elif file_type == "ppt":
                return self._generate_ppt(content, title, suggested_filename, output_dir)
            elif file_type == "png":
                return await self._generate_png(content, title, suggested_filename, output_dir)
            else:
                return json.dumps({"error": f"Unsupported file_type: {file_type}"})
        except Exception as e:
            return json.dumps({"error": f"File generation failed: {e}"})

    def _make_filename(self, suggested: str, file_type: str) -> str:
        if suggested:
            base = re.sub(r"[^\w\-.]", "_", suggested).strip("_")
            if not base:
                base = f"file-{uuid.uuid4().hex[:8]}"
            ext_map = {"pdf": ".pdf", "doc": ".docx", "ppt": ".pptx", "png": ".png"}
            expected_ext = ext_map.get(file_type, "")
            if expected_ext and not base.lower().endswith(expected_ext):
                base = base + expected_ext
            return base
        return f"file-{uuid.uuid4().hex[:8]}.{file_type}"

    def _downloadable_response(self, filepath: Path, content_type: str) -> str:
        size = filepath.stat().st_size
        filename = filepath.name
        # ``download_url`` points at /api/download (forces an attachment
        # via Content-Disposition so clicking the link always saves the
        # file, regardless of MIME type). ``preview_url`` points at
        # /api/files (serves the file inline with its proper MIME type)
        # so the canvas can embed PDFs in an iframe and images in <img>.
        return json.dumps(
            {
                "downloadable": {
                    "filename": filename,
                    "content_type": content_type,
                    "size": size,
                    "download_url": f"/api/download/{filename}",
                    "preview_url": f"/api/files/{filename}",
                }
            }
        )

    # ─────────────────────────── PDF ───────────────────────────

    # Transliteration map for characters the latin-1 core fonts (Helvetica)
    # cannot render. Used only when no Unicode TTF font is available.
    _LATIN1_FALLBACK = {
        "•": "-", "◦": "-", "▪": "-", "‣": "-", "·": "-", "●": "-",
        "“": '"', "”": '"', "„": '"', "‘": "'", "’": "'", "‚": "'",
        "–": "-", "—": "-", "―": "-", "−": "-",
        "…": "...", "→": "->", "←": "<-", "⇒": "=>", "⇐": "<=",
        "™": "(TM)", "®": "(R)", "©": "(C)", "°": " deg",
        "≤": "<=", "≥": ">=", "≠": "!=", "×": "x", "÷": "/",
        " ": " ", "​": "", "﻿": "", "\t": "    ",
    }

    def _register_pdf_font(self, pdf) -> tuple[str, bool]:
        """Register a Unicode TTF font if one is available on the system.

        Returns ``(family, unicode_ok)``. When a TTF is found the document
        can render arbitrary Unicode (bullets, smart quotes, accents …).
        Otherwise we fall back to the built-in Helvetica core font, which is
        limited to latin-1 and requires the caller to transliterate text.
        """
        candidates = [
            (
                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                "/usr/share/fonts/truetype/dejavu/DejaVuSans-Bold.ttf",
            ),
            (
                "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
                "/usr/share/fonts/truetype/liberation/LiberationSans-Bold.ttf",
            ),
            (
                "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
                "/usr/share/fonts/truetype/freefont/FreeSansBold.ttf",
            ),
        ]
        for regular, bold in candidates:
            if not os.path.exists(regular):
                continue
            try:
                pdf.add_font("DocFont", "", regular)
                pdf.add_font("DocFont", "B", bold if os.path.exists(bold) else regular)
                return "DocFont", True
            except Exception as e:  # pragma: no cover - font load is environment-specific
                logger.warning("Failed to register Unicode font %s: %s", regular, e)
        return "Helvetica", False

    def _generate_pdf(
        self, content: str, title: str, suggested_filename: str, output_dir: Path
    ) -> str:
        try:
            from fpdf import FPDF
        except ImportError:
            return json.dumps({"error": "fpdf2 is not installed. Run: pip install fpdf2"})

        filename = self._make_filename(suggested_filename, "pdf")
        filepath = output_dir / filename

        pdf = FPDF()
        pdf.set_auto_page_break(auto=True, margin=15)
        pdf.add_page()

        font_family, unicode_ok = self._register_pdf_font(pdf)
        bullet = "• " if unicode_ok else "- "

        def sanitize(text: str) -> str:
            if unicode_ok:
                return text
            for bad, good in self._LATIN1_FALLBACK.items():
                text = text.replace(bad, good)
            # Drop anything still outside latin-1 so fpdf never raises.
            return text.encode("latin-1", "replace").decode("latin-1")

        def write(text: str, size: int, bold: bool = False, indent: float = 0.0) -> None:
            """Render a block of text safely.

            ``wrapmode="CHAR"`` lets fpdf break overly long unbroken tokens
            (URLs, hashes …) at the character level instead of raising
            "Not enough horizontal space to render a single character".
            """
            pdf.set_font(font_family, "B" if bold else "", size)
            text = sanitize(text)
            if not text:
                pdf.ln(size * 0.35)
                return
            width = max(pdf.epw - indent, 1.0)
            if indent:
                pdf.set_x(pdf.l_margin + indent)
            pdf.multi_cell(
                width,
                size * 0.5,
                text,
                align="L",
                new_x="LMARGIN",
                new_y="NEXT",
                wrapmode="CHAR",
            )

        if title:
            pdf.set_font(font_family, "B", 16)
            pdf.multi_cell(
                0, 10, sanitize(title), align="C",
                new_x="LMARGIN", new_y="NEXT", wrapmode="CHAR",
            )
            pdf.ln(5)

        for raw_line in content.splitlines():
            line = raw_line.rstrip()
            if not line:
                pdf.ln(4)
            elif line.startswith("### "):
                write(line[4:].strip(), 13, bold=True)
            elif line.startswith("## "):
                write(line[3:].strip(), 14, bold=True)
            elif line.startswith("# "):
                write(line[2:].strip(), 16, bold=True)
            elif line.startswith("- ") or line.startswith("* "):
                write(bullet + line[2:].strip(), 12, indent=5)
            elif re.match(r"^\d+\.\s", line):
                write(line.strip(), 12, indent=5)
            else:
                write(line, 12)

        pdf.output(str(filepath))
        return self._downloadable_response(filepath, "application/pdf")

    # ─────────────────────────── DOCX ───────────────────────────

    def _generate_doc(
        self, content: str, title: str, suggested_filename: str, output_dir: Path
    ) -> str:
        try:
            from docx import Document
            from docx.enum.text import WD_ALIGN_PARAGRAPH
            from docx.shared import Inches, Pt
        except ImportError:
            return json.dumps(
                {"error": "python-docx is not installed. Run: pip install python-docx"}
            )

        filename = self._make_filename(suggested_filename, "doc")
        filepath = output_dir / filename

        doc = Document()

        if title:
            heading = doc.add_heading(title, level=0)
            heading.alignment = WD_ALIGN_PARAGRAPH.CENTER

        lines = content.splitlines()
        for line in lines:
            stripped = line.rstrip()
            if not stripped:
                continue

            if stripped.startswith("### "):
                doc.add_heading(stripped[4:].strip(), level=3)
            elif stripped.startswith("## "):
                doc.add_heading(stripped[3:].strip(), level=2)
            elif stripped.startswith("# "):
                doc.add_heading(stripped[2:].strip(), level=1)
            elif stripped.startswith("- ") or stripped.startswith("* "):
                p = doc.add_paragraph(stripped[2:].strip(), style="List Bullet")
            elif re.match(r"^\d+\.\s", stripped):
                p = doc.add_paragraph(stripped.strip(), style="List Number")
            else:
                doc.add_paragraph(stripped)

        doc.save(str(filepath))
        return self._downloadable_response(
            filepath, "application/vnd.openxmlformats-officedocument.wordprocessingml.document"
        )

    # ─────────────────────────── PPTX ───────────────────────────

    def _generate_ppt(
        self, content: str, title: str, suggested_filename: str, output_dir: Path
    ) -> str:
        try:
            from pptx import Presentation
            from pptx.util import Inches, Pt
        except ImportError:
            return json.dumps(
                {"error": "python-pptx is not installed. Run: pip install python-pptx"}
            )

        filename = self._make_filename(suggested_filename, "ppt")
        filepath = output_dir / filename

        prs = Presentation()
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)

        slide_blocks = re.split(r"\n?---+\n?", content)
        if len(slide_blocks) == 1:
            lines = content.splitlines()
            slide_blocks = []
            current = []
            for line in lines:
                if re.match(r"^#\s", line) and current:
                    slide_blocks.append("\n".join(current))
                    current = [line]
                else:
                    current.append(line)
            if current:
                slide_blocks.append("\n".join(current))

        for block in slide_blocks:
            block = block.strip()
            if not block:
                continue

            slide_layout = prs.slide_layouts[1]
            slide = prs.slides.add_slide(slide_layout)
            shapes = slide.shapes

            title_text = ""
            body_lines = []
            lines = block.splitlines()
            title_found = False
            for line in lines:
                if not title_found and re.match(r"^#+\s", line):
                    title_text = re.sub(r"^#+\s", "", line).strip()
                    title_found = True
                else:
                    body_lines.append(line)

            if not title_text:
                title_text = title or "Slide"

            if shapes.title:
                shapes.title.text = title_text

            body_text = "\n".join(body_lines).strip()
            if body_text and len(shapes.placeholders) > 1:
                tf = shapes.placeholders[1].text_frame
                tf.clear()
                for bline in body_lines:
                    bline = bline.rstrip()
                    if not bline:
                        continue
                    if bline.startswith("- ") or bline.startswith("* "):
                        p = tf.add_paragraph()
                        p.text = bline[2:].strip()
                        p.level = 0
                    elif re.match(r"^\d+\.\s", bline):
                        p = tf.add_paragraph()
                        p.text = bline.strip()
                        p.level = 0
                    else:
                        p = tf.add_paragraph()
                        p.text = bline
                        p.level = 0

        prs.save(str(filepath))
        return self._downloadable_response(
            filepath, "application/vnd.openxmlformats-officedocument.presentationml.presentation"
        )

    # ─────────────────────────── PNG ───────────────────────────

    async def _generate_png(
        self, content: str, title: str, suggested_filename: str, output_dir: Path
    ) -> str:
        try:
            from PIL import Image, ImageDraw, ImageFont
        except ImportError:
            return json.dumps({"error": "Pillow is not installed. Run: pip install Pillow"})

        filename = self._make_filename(suggested_filename, "png")
        filepath = output_dir / filename

        # 1) Try base64 decode first (fast, sync)
        image = self._try_base64_to_image(content)
        if image:
            image.save(str(filepath), "PNG")
            return self._downloadable_response(filepath, "image/png")

        # 2) Try AI text-to-image generation
        #    Offload the heavy model loading + inference to a real OS thread
        #    so the asyncio event loop stays responsive.
        try:
            cache_dir = self.config.get("model_cache_dir")
            await asyncio.to_thread(_ensure_diffusion_loaded, cache_dir)
            image = await asyncio.to_thread(_run_inference, content, title)
            image.save(str(filepath), "PNG")
            return self._downloadable_response(filepath, "image/png")
        except RuntimeError as e:
            logger.warning("Diffusion model unavailable (%s). Falling back to text rendering.", e)
        except Exception as e:
            logger.warning("Diffusion generation failed (%s). Falling back to text rendering.", e)

        # 3) Fall back: render text as image (fast, sync)
        img_width, img_height = 1200, 800
        img = Image.new("RGB", (img_width, img_height), "#f7f6f3")
        draw = ImageDraw.Draw(img)

        font = self._get_font(20)
        title_font = self._get_font(28)

        y = 30
        if title:
            draw.text((30, y), title, fill="#2c2c2c", font=title_font)
            y += 50

        margin = 30
        max_width = img_width - 2 * margin
        line_height = 28

        lines = content.splitlines()
        for line in lines:
            if y > img_height - margin:
                break
            stripped = line.rstrip()
            if not stripped:
                y += line_height // 2
                continue

            words = stripped.split(" ")
            current_line = ""
            for word in words:
                test = current_line + " " + word if current_line else word
                bbox = draw.textbbox((0, 0), test, font=font)
                if bbox[2] <= max_width:
                    current_line = test
                else:
                    draw.text((margin, y), current_line, fill="#2c2c2c", font=font)
                    y += line_height
                    current_line = word
            if current_line:
                draw.text((margin, y), current_line, fill="#2c2c2c", font=font)
                y += line_height

        img.save(str(filepath), "PNG")
        return self._downloadable_response(filepath, "image/png")

    def _try_base64_to_image(self, content: str):
        try:
            import io

            from PIL import Image

            cleaned = re.sub(r"^data:image/[^;]+;base64,", "", content.strip())
            cleaned = re.sub(r"\s", "", cleaned)
            decoded = base64.b64decode(cleaned)
            return Image.open(io.BytesIO(decoded))
        except Exception:
            return None

    def _get_font(self, size: int):
        try:
            from PIL import ImageFont

            candidates = [
                "/usr/share/fonts/truetype/dejavu/DejaVuSans.ttf",
                "/usr/share/fonts/truetype/liberation/LiberationSans-Regular.ttf",
                "/usr/share/fonts/truetype/freefont/FreeSans.ttf",
                "/System/Library/Fonts/Helvetica.ttc",
                "C:\\Windows\\Fonts\\arial.ttf",
            ]
            for path in candidates:
                if os.path.exists(path):
                    return ImageFont.truetype(path, size)
        except Exception:
            pass
        from PIL import ImageFont

        return ImageFont.load_default()
